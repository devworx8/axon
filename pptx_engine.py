"""Axon PPTX Engine — generate polished PowerPoint presentations.

Supports:
  - Structured slide decks from JSON slide data
  - AI-prompted decks (pass prompt, get back a deck)
  - Dark (default) and light themes
  - Output to ~/Downloads or a custom path
"""
from __future__ import annotations

import json
import os
import re
from dataclasses import dataclass, field
from datetime import date
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu


# ── Colour palettes ──────────────────────────────────────────────────────────

THEMES: dict[str, dict[str, Any]] = {
    "dark": {
        "bg":           RGBColor(0x02, 0x06, 0x17),   # slate-950
        "slide_bg":     RGBColor(0x0F, 0x17, 0x2A),   # slate-900
        "accent":       RGBColor(0x38, 0xBD, 0xF8),   # sky-400
        "accent2":      RGBColor(0xA7, 0x8B, 0xFA),   # violet-400
        "title_text":   RGBColor(0xF1, 0xF5, 0xF9),   # slate-100
        "body_text":    RGBColor(0xCB, 0xD5, 0xE1),   # slate-300
        "muted_text":   RGBColor(0x64, 0x74, 0x8B),   # slate-500
        "bullet_dot":   RGBColor(0x38, 0xBD, 0xF8),   # sky-400
        "divider":      RGBColor(0x1E, 0x29, 0x3B),   # slate-800
        "tag_bg":       RGBColor(0x0E, 0xA5, 0xE9),   # sky-500
    },
    "light": {
        "bg":           RGBColor(0xF8, 0xFA, 0xFC),
        "slide_bg":     RGBColor(0xFF, 0xFF, 0xFF),
        "accent":       RGBColor(0x0E, 0xA5, 0xE9),
        "accent2":      RGBColor(0x7C, 0x3A, 0xED),
        "title_text":   RGBColor(0x0F, 0x17, 0x2A),
        "body_text":    RGBColor(0x1E, 0x29, 0x3B),
        "muted_text":   RGBColor(0x64, 0x74, 0x8B),
        "bullet_dot":   RGBColor(0x0E, 0xA5, 0xE9),
        "divider":      RGBColor(0xE2, 0xE8, 0xF0),
        "tag_bg":       RGBColor(0x0E, 0xA5, 0xE9),
    },
}


# ── Slide data model ─────────────────────────────────────────────────────────

@dataclass
class SlideData:
    """One slide in the deck."""
    type: str          # "title" | "section" | "bullets" | "two_col" | "closing"
    title: str = ""
    subtitle: str = ""
    bullets: list[str] = field(default_factory=list)
    left: list[str] = field(default_factory=list)
    right: list[str] = field(default_factory=list)
    left_title: str = ""
    right_title: str = ""
    tag: str = ""      # small badge label (e.g. "Agenda", "Finance")
    notes: str = ""    # speaker notes


@dataclass
class DeckSpec:
    """Full deck specification."""
    title: str
    subtitle: str = ""
    author: str = ""
    date: str = ""
    theme: str = "dark"
    slides: list[SlideData] = field(default_factory=list)
    output_path: str = ""   # empty = auto ~/Downloads


# ── Helper: set solid fill on a shape ────────────────────────────────────────

def _solid(shape, colour: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = colour


def _text_run(para, text: str, *, size: int, bold: bool = False,
               colour: RGBColor | None = None, italic: bool = False) -> None:
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if colour:
        run.font.color.rgb = colour


def _add_textbox(slide, left, top, width, height,
                 text: str, size: int, colour: RGBColor,
                 bold: bool = False, align=PP_ALIGN.LEFT, italic: bool = False):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = align
    _text_run(para, text, size=size, bold=bold, colour=colour, italic=italic)
    return txb


# ── Slide builders ────────────────────────────────────────────────────────────

SW = Inches(13.33)   # widescreen width
SH = Inches(7.5)     # widescreen height
M  = Inches(0.55)    # margin


def _build_title_slide(slide, sd: SlideData, C: dict) -> None:
    """Full-bleed title card with accent bar."""
    _solid(slide.background, C["slide_bg"])

    # Accent bar left edge
    bar = slide.shapes.add_shape(1, 0, Inches(1.8), Inches(0.06), SH - Inches(1.8))
    _solid(bar, C["accent"])
    bar.line.fill.background()

    # Tag badge
    if sd.tag:
        badge_w = Inches(1.8)
        badge = slide.shapes.add_shape(1, M, Inches(1.0), badge_w, Inches(0.35))
        _solid(badge, C["tag_bg"])
        badge.line.fill.background()
        _add_textbox(slide, M + Inches(0.12), Inches(1.0), badge_w, Inches(0.35),
                     sd.tag.upper(), 9, C["slide_bg"], bold=True)

    # Title
    _add_textbox(slide, Inches(1.2), Inches(2.0), Inches(9.0), Inches(1.8),
                 sd.title, 40, C["title_text"], bold=True)

    # Subtitle
    if sd.subtitle:
        _add_textbox(slide, Inches(1.2), Inches(3.9), Inches(9.0), Inches(1.2),
                     sd.subtitle, 20, C["body_text"], italic=True)

    # Bottom meta row
    meta = []
    if sd.notes:
        meta.append(sd.notes)
    _add_textbox(slide, Inches(1.2), SH - Inches(0.9), Inches(9.0), Inches(0.5),
                 " · ".join(meta) if meta else "", 11, C["muted_text"])


def _build_section_slide(slide, sd: SlideData, C: dict) -> None:
    """Section divider — accent background strip."""
    _solid(slide.background, C["bg"])

    strip = slide.shapes.add_shape(1, 0, Inches(2.5), SW, Inches(2.5))
    _solid(strip, C["slide_bg"])
    strip.line.fill.background()

    if sd.tag:
        _add_textbox(slide, M, Inches(2.65), Inches(4), Inches(0.4),
                     sd.tag.upper(), 10, C["accent"], bold=True)

    _add_textbox(slide, M, Inches(3.1), Inches(10), Inches(1.4),
                 sd.title, 36, C["title_text"], bold=True)

    if sd.subtitle:
        _add_textbox(slide, M, Inches(4.55), Inches(9), Inches(0.8),
                     sd.subtitle, 17, C["body_text"])


def _build_bullets_slide(slide, sd: SlideData, C: dict) -> None:
    """Standard title + bullet list slide."""
    _solid(slide.background, C["slide_bg"])

    # Accent top bar
    bar = slide.shapes.add_shape(1, 0, 0, SW, Inches(0.06))
    _solid(bar, C["accent"])
    bar.line.fill.background()

    # Tag
    if sd.tag:
        _add_textbox(slide, M, Inches(0.3), Inches(3), Inches(0.35),
                     sd.tag.upper(), 9, C["accent"], bold=True)

    # Title
    _add_textbox(slide, M, Inches(0.7), Inches(11.5), Inches(0.9),
                 sd.title, 26, C["title_text"], bold=True)

    # Divider line
    div = slide.shapes.add_shape(1, M, Inches(1.65), Inches(11.2), Inches(0.025))
    _solid(div, C["divider"])
    div.line.fill.background()

    # Bullets
    top = Inches(1.85)
    row_h = Inches(0.52)
    for i, bullet in enumerate(sd.bullets):
        y = top + i * row_h
        # Dot
        dot = slide.shapes.add_shape(9, M, y + Inches(0.13), Inches(0.12), Inches(0.12))
        _solid(dot, C["bullet_dot"])
        dot.line.fill.background()
        # Text
        _add_textbox(slide, M + Inches(0.28), y, Inches(11.8), row_h,
                     bullet, 16, C["body_text"])


def _build_two_col_slide(slide, sd: SlideData, C: dict) -> None:
    """Two-column layout."""
    _solid(slide.background, C["slide_bg"])

    bar = slide.shapes.add_shape(1, 0, 0, SW, Inches(0.06))
    _solid(bar, C["accent"])
    bar.line.fill.background()

    if sd.tag:
        _add_textbox(slide, M, Inches(0.3), Inches(3), Inches(0.35),
                     sd.tag.upper(), 9, C["accent"], bold=True)

    _add_textbox(slide, M, Inches(0.7), Inches(11.5), Inches(0.9),
                 sd.title, 26, C["title_text"], bold=True)

    div = slide.shapes.add_shape(1, M, Inches(1.65), Inches(11.2), Inches(0.025))
    _solid(div, C["divider"])
    div.line.fill.background()

    col_w = Inches(5.8)
    col_gap = Inches(0.5)
    left_x = M
    right_x = M + col_w + col_gap

    # Column titles
    if sd.left_title:
        _add_textbox(slide, left_x, Inches(1.9), col_w, Inches(0.45),
                     sd.left_title, 14, C["accent"], bold=True)
    if sd.right_title:
        _add_textbox(slide, right_x, Inches(1.9), col_w, Inches(0.45),
                     sd.right_title, 14, C["accent2"], bold=True)

    # Vertical divider
    vdiv = slide.shapes.add_shape(1, M + col_w + Inches(0.22), Inches(1.85),
                                   Inches(0.025), Inches(5.2))
    _solid(vdiv, C["divider"])
    vdiv.line.fill.background()

    row_h = Inches(0.5)
    top = Inches(2.45)

    for i, item in enumerate(sd.left):
        y = top + i * row_h
        dot = slide.shapes.add_shape(9, left_x, y + Inches(0.13), Inches(0.1), Inches(0.1))
        _solid(dot, C["bullet_dot"])
        dot.line.fill.background()
        _add_textbox(slide, left_x + Inches(0.22), y, col_w - Inches(0.25), row_h,
                     item, 15, C["body_text"])

    for i, item in enumerate(sd.right):
        y = top + i * row_h
        dot = slide.shapes.add_shape(9, right_x, y + Inches(0.13), Inches(0.1), Inches(0.1))
        _solid(dot, C["accent2"])
        dot.line.fill.background()
        _add_textbox(slide, right_x + Inches(0.22), y, col_w - Inches(0.25), row_h,
                     item, 15, C["body_text"])


def _build_closing_slide(slide, sd: SlideData, C: dict) -> None:
    """Thank-you / closing slide."""
    _solid(slide.background, C["bg"])

    # Full-width accent strip
    strip = slide.shapes.add_shape(1, 0, Inches(2.8), SW, Inches(0.08))
    _solid(strip, C["accent"])
    strip.line.fill.background()

    _add_textbox(slide, 0, Inches(2.2), SW, Inches(1.2),
                 sd.title, 42, C["title_text"], bold=True, align=PP_ALIGN.CENTER)

    if sd.subtitle:
        _add_textbox(slide, 0, Inches(3.5), SW, Inches(0.8),
                     sd.subtitle, 20, C["body_text"], align=PP_ALIGN.CENTER, italic=True)

    if sd.notes:
        _add_textbox(slide, 0, SH - Inches(1.0), SW, Inches(0.5),
                     sd.notes, 12, C["muted_text"], align=PP_ALIGN.CENTER)


# ── Slide dispatcher ──────────────────────────────────────────────────────────

_BUILDERS = {
    "title":    _build_title_slide,
    "section":  _build_section_slide,
    "bullets":  _build_bullets_slide,
    "two_col":  _build_two_col_slide,
    "closing":  _build_closing_slide,
}


def build_deck(spec: DeckSpec) -> Path:
    """Build the PPTX from a DeckSpec and return the saved file path."""
    C = THEMES.get(spec.theme, THEMES["dark"])

    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    blank_layout = prs.slide_layouts[6]   # truly blank

    for sd in spec.slides:
        slide = prs.slides.add_slide(blank_layout)
        builder = _BUILDERS.get(sd.type, _build_bullets_slide)
        builder(slide, sd, C)
        if sd.notes:
            slide.notes_slide.notes_text_frame.text = sd.notes

    # Determine output path
    if spec.output_path:
        out = Path(spec.output_path).expanduser()
    else:
        downloads = Path.home() / "Downloads"
        downloads.mkdir(exist_ok=True)
        safe_title = re.sub(r"[^\w\s-]", "", spec.title).strip().replace(" ", "_")
        out = downloads / f"{safe_title}_{date.today().isoformat()}.pptx"

    prs.save(str(out))
    return out


# ── JSON → DeckSpec parser ────────────────────────────────────────────────────

def deck_from_dict(data: dict[str, Any]) -> DeckSpec:
    slides = []
    for s in data.get("slides", []):
        slides.append(SlideData(
            type=s.get("type", "bullets"),
            title=s.get("title", ""),
            subtitle=s.get("subtitle", ""),
            bullets=s.get("bullets", []),
            left=s.get("left", []),
            right=s.get("right", []),
            left_title=s.get("left_title", ""),
            right_title=s.get("right_title", ""),
            tag=s.get("tag", ""),
            notes=s.get("notes", ""),
        ))
    return DeckSpec(
        title=data.get("title", "Presentation"),
        subtitle=data.get("subtitle", ""),
        author=data.get("author", ""),
        date=data.get("date", date.today().isoformat()),
        theme=data.get("theme", "dark"),
        slides=slides,
        output_path=data.get("output_path", ""),
    )


# ── AI prompt → slide JSON ────────────────────────────────────────────────────

SYSTEM_PROMPT = """You are a presentation designer. Given a topic and context,
return a JSON deck specification. Use this exact schema:

{
  "title": "Deck title",
  "subtitle": "Optional subtitle",
  "theme": "dark",
  "slides": [
    {
      "type": "title",
      "title": "...",
      "subtitle": "...",
      "tag": "Optional badge",
      "notes": "Speaker notes"
    },
    {
      "type": "bullets",
      "title": "Slide heading",
      "tag": "Optional",
      "bullets": ["Point 1", "Point 2", "Point 3"],
      "notes": "Speaker notes"
    },
    {
      "type": "two_col",
      "title": "Comparison heading",
      "left_title": "Left column heading",
      "right_title": "Right column heading",
      "left": ["Item 1", "Item 2"],
      "right": ["Item A", "Item B"],
      "notes": "Speaker notes"
    },
    {
      "type": "section",
      "title": "Section name",
      "subtitle": "Section intro",
      "tag": "Part 1"
    },
    {
      "type": "closing",
      "title": "Thank You",
      "subtitle": "Closing message",
      "notes": "Contact info or next steps"
    }
  ]
}

Rules:
- Always start with a "title" slide and end with a "closing" slide
- 5–12 slides total
- Max 6 bullets per slide
- Keep bullet text concise (under 10 words each)
- Use "two_col" for comparisons or before/after content
- Use "section" as dividers between major topics
- Return ONLY valid JSON, no markdown fences
"""


def prompt_to_deck_json(prompt: str, context: str = "", model_fn=None) -> dict:
    """Call the AI model to generate slide JSON from a prompt.
    model_fn: callable(system, user) -> str  (injected by server)
    """
    if model_fn is None:
        raise RuntimeError("model_fn required for AI deck generation")

    user_msg = f"Topic: {prompt}"
    if context:
        user_msg += f"\n\nContext:\n{context}"

    raw = model_fn(SYSTEM_PROMPT, user_msg)

    # Strip markdown fences defensively
    raw = re.sub(r"^```(?:json)?\s*", "", raw.strip(), flags=re.IGNORECASE)
    raw = re.sub(r"\s*```$", "", raw.strip())

    return json.loads(raw)
