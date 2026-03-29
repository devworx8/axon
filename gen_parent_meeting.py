#!/usr/bin/env python3
"""Generate Parent Meeting PPTX for Young Eagles."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

out = os.path.expanduser("~/Desktop/Parent_Meeting_Young_Eagles.pptx")
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

bg_dark = RGBColor(0x1B, 0x2A, 0x4A)
accent = RGBColor(0x00, 0x96, 0xD6)
white = RGBColor(0xFF, 0xFF, 0xFF)
light_gray = RGBColor(0xCC, 0xCC, 0xCC)
gold = RGBColor(0xFF, 0xB8, 0x00)

def set_bg(slide, color=bg_dark):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=white, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf

def add_bullet_slide(title_text, bullets, subtitle_text=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide)
    # accent bar
    slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.08))
    bar = slide.shapes[-1]
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    # title
    add_text_box(slide, 0.8, 0.4, 11, 1, title_text, 36, accent, True)
    if subtitle_text:
        add_text_box(slide, 0.8, 1.1, 11, 0.6, subtitle_text, 16, light_gray)
    # bullets
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.font.size = Pt(20)
        p.font.color.rgb = white
        p.space_after = Pt(12)
        p.level = 0
    return slide

# ━━━━━━━━━ SLIDE 1: TITLE ━━━━━━━━━
s1 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s1)
add_text_box(s1, 1, 1.5, 11, 1.5, "YOUNG EAGLES", 52, accent, True, PP_ALIGN.CENTER)
add_text_box(s1, 1, 3.0, 11, 1.2, "Parent Meeting", 40, white, True, PP_ALIGN.CENTER)
add_text_box(s1, 1, 4.3, 11, 0.8, "27 March 2026", 22, light_gray, False, PP_ALIGN.CENTER)
add_text_box(s1, 1, 5.3, 11, 0.6, "Building Strong Futures Together", 18, gold, False, PP_ALIGN.CENTER)

# ━━━━━━━━━ SLIDE 2: AGENDA ━━━━━━━━━
add_bullet_slide("Meeting Agenda", [
    "1.  Uniform",
    "2.  Excursion",
    "3.  Conduct",
    "4.  Fundraising — Aircon Project",
    "5.  Projects",
    "6.  Academics — Reports & Participation",
    "7.  Year End Function — 31 October 2026 (Graduation Ceremony)",
])

# ━━━━━━━━━ SLIDE 3: GROUND RULES ━━━━━━━━━
add_bullet_slide("Ground Rules", [
    "✋  Please raise your hand before speaking",
    "📋  Questions must relate to the current agenda item only",
    "🚫  Topics outside of today's agenda will not be addressed",
    "⏱️  Keep questions brief — we have a lot to cover",
    "🤝  Be respectful of everyone's time and opinions",
    "📵  Please silence your phones",
    "📝  A summary will be shared after the meeting",
], "Let's keep this meeting focused and productive")

# ━━━━━━━━━ SLIDE 4: UNIFORM ━━━━━━━━━
add_bullet_slide("1. Uniform", [
    "All learners must wear the full school uniform daily",
    "Uniform must be neat, clean, and properly fitted",
    "No casual wear allowed on school days unless communicated",
    "Lost items — check the lost & found before purchasing replacements",
    "Winter uniform transition dates will be communicated",
])

# ━━━━━━━━━ SLIDE 5: EXCURSION ━━━━━━━━━
add_bullet_slide("2. Excursion", [
    "Upcoming excursion planned for this term",
    "Permission slips must be signed and returned on time",
    "Outstanding payments must be settled before the trip",
    "Safety briefing will be provided to all learners beforehand",
    "Details of destination and date to follow via communication",
])

# ━━━━━━━━━ SLIDE 6: CONDUCT — LEARNERS ━━━━━━━━━
add_bullet_slide("3. Conduct — Learner Expectations", [
    "All learners must uphold the school's Code of Conduct",
    "Bullying, disrespect, and disruption will not be tolerated",
    "Parents will be contacted for repeated behavioural issues",
    "Positive behaviour is recognised and rewarded",
    "Please reinforce good conduct and discipline at home",
])

# ━━━━━━━━━ SLIDE 6b: CONDUCT — PARENTS ━━━━━━━━━
add_bullet_slide("3. Conduct — Parent Expectations", [
    "Treat all staff, teachers, and other parents with respect",
    "Do not confront or address teachers in front of learners",
    "Follow the proper channels — class teacher → HOD → principal",
    "No parent may enter classrooms without prior arrangement",
    "Abusive language or threatening behaviour will not be tolerated",
    "Support the school's disciplinary decisions — work with us, not against us",
    "Model the behaviour you want your child to display",
])

# ━━━━━━━━━ SLIDE 7: FUNDRAISING ━━━━━━━━━
add_bullet_slide("4. Fundraising — Aircon Project", [
    "Goal: Install air conditioning in classrooms",
    "Fundraising events and contributions are welcome",
    "Every contribution counts — no amount is too small",
    "Updates on funds raised will be shared regularly",
    "Volunteers needed to help organise fundraising activities",
])

# ━━━━━━━━━ SLIDE 8: PROJECTS ━━━━━━━━━
add_bullet_slide("5. Projects", [
    "Upcoming projects for this term have been assigned",
    "Parents are encouraged to supervise but not do the work",
    "Deadlines are firm — late submissions affect marks",
    "Materials lists will be shared in advance",
    "Reach out to the class teacher for any project-related concerns",
])

# ━━━━━━━━━ SLIDE 9: ACADEMICS ━━━━━━━━━
add_bullet_slide("6. Academics", [
    "📊  Term reports will be issued — please review them carefully",
    "Participation in class is part of the assessment criteria",
    "Homework must be completed daily — please check at home",
    "Extra support is available — speak to the teacher if needed",
    "Parent-teacher consultations can be arranged on request",
], "Reports & Participation")

# ━━━━━━━━━ SLIDE 10: YEAR END FUNCTION ━━━━━━━━━
add_bullet_slide("7. Year End Function — Graduation Ceremony", [
    "📅  Date: 31 October 2026",
    "🎓  Graduation Ceremony for qualifying learners",
    "👗  Dress code and programme details to follow",
    "📸  Photo opportunities will be arranged",
    "🎉  Families are welcome to attend and celebrate",
    "Planning committee volunteers are welcome",
], "Save the date!")

# ━━━━━━━━━ SLIDE 11: THANK YOU ━━━━━━━━━
s_end = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s_end)
add_text_box(s_end, 1, 2.0, 11, 1.2, "Thank You", 48, accent, True, PP_ALIGN.CENTER)
add_text_box(s_end, 1, 3.5, 11, 1, "Together we build strong futures for our children",
             24, white, False, PP_ALIGN.CENTER)
add_text_box(s_end, 1, 4.8, 11, 0.8, "Young Eagles — Soaring Higher",
             20, gold, False, PP_ALIGN.CENTER)
add_text_box(s_end, 1, 5.8, 11, 0.6, "Questions? Please contact the school office",
             16, light_gray, False, PP_ALIGN.CENTER)

prs.save(out)
print(f"✅ Saved: {out}")
print(f"   {len(prs.slides)} slides generated")
